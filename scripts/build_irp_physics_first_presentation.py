from __future__ import annotations

import math
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib import cm
from matplotlib.colors import LinearSegmentedColormap, Normalize
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
BOUND_PATH = ROOT / "extraction_outputs" / "bound_outcomes.csv"
PREDICTIONS_PATH = ROOT / "ml" / "physics_structured_surrogate" / "tables" / "predictions_with_trust_flags.csv"
FIG_DIR = ROOT / "figures" / "presentation"
OUT_PATH = ROOT / "deliverables" / "IRP_progress_physics_first_10_slides_2026-08-11.pptx"

BLUE = "#1f5aa6"
LIGHT_BLUE = "#8bb8f0"
DARK = "#102033"
GREY = "#d9dde4"
MID_GREY = "#6d7887"
PALE = "#f4f7fb"
ACCENT = "#c86a2b"
SPIN_COLORS = {"x": "#d95f02", "y": "#1b9e77", "z": "#7570b3", "mz": "#e7298a"}


def parse_numeric_code(series: pd.Series, pattern: str, scale: float = 1.0) -> pd.Series:
    extracted = series.fillna("").astype(str).str.extract(pattern)[0]
    return pd.to_numeric(extracted, errors="coerce") / scale


def load_bound_frame() -> pd.DataFrame:
    frame = pd.read_csv(BOUND_PATH, low_memory=False)
    frame["mass_log10_kg"] = parse_numeric_code(frame["mass_code"], r"A(\d{4})", 100.0)
    frame["periapsis_Rm"] = parse_numeric_code(frame["periapsis_code"], r"r(\d+)", 10.0)
    frame["v_inf_kms"] = parse_numeric_code(frame["velocity_code"], r"v(\d+)", 10.0)
    frame["spin_period_hr"] = parse_numeric_code(frame["spin_code"], r"s(\d{3})", 10.0)
    frame["spin_label"] = "none"
    spin_code = frame["spin_code"].fillna("").astype(str)
    frame.loc[spin_code.str.contains("mz"), "spin_label"] = "mz"
    frame.loc[spin_code.str.contains("x"), "spin_label"] = "x"
    frame.loc[spin_code.str.contains("y"), "spin_label"] = "y"
    frame.loc[spin_code.str.contains("z") & ~spin_code.str.contains("mz"), "spin_label"] = "z"
    frame["has_spin"] = spin_code.ne("")
    frame["target_mass_kg"] = np.power(10.0, frame["mass_log10_kg"])
    frame["parent_mass_kg"] = frame["target_mass_kg"]
    frame["parent_radius_km"] = np.cbrt((3.0 * frame["parent_mass_kg"]) / (4.0 * np.pi * 2700.0)) / 1000.0
    return frame


def dedupe_physical_runs(frame: pd.DataFrame) -> pd.DataFrame:
    group_cols = [
        "physical_file",
        "mass_log10_kg",
        "periapsis_Rm",
        "v_inf_kms",
        "spin_label",
        "spin_period_hr",
        "has_spin",
        "resolution_code",
    ]
    return (
        frame.groupby(group_cols, dropna=False)
        .agg(
            bound_mass_fraction=("bound_mass_fraction", "mean"),
            parent_radius_km=("parent_radius_km", "mean"),
            parent_mass_kg=("parent_mass_kg", "mean"),
        )
        .reset_index()
    )


def load_prediction_summary() -> pd.DataFrame:
    pred = pd.read_csv(PREDICTIONS_PATH, low_memory=False)
    return (
        pred.groupby(["physical_file", "mass_log10_kg", "periapsis_Rm", "v_inf_kms"], as_index=False)
        .agg(
            actual_bmf=("bound_mass_fraction", "mean"),
            oof_prediction=("predicted", "mean"),
            abs_error=("residual", lambda s: s.abs().mean()),
        )
    )


def support_note(items: list[tuple[str, int]]) -> str:
    return " | ".join(f"{label}: n={count}" for label, count in items)


def add_axis_header(ax: plt.Axes, title: str, note: str) -> None:
    ax.set_title(title, loc="left", fontsize=15, fontweight="bold", color=DARK, pad=14)
    ax.text(
        0.0,
        1.02,
        note,
        transform=ax.transAxes,
        fontsize=9,
        color=MID_GREY,
        va="bottom",
    )


def style_axes(ax: plt.Axes) -> None:
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#95a1af")
    ax.spines["bottom"].set_color("#95a1af")
    ax.grid(True, axis="y", color="#e6ebf2", linewidth=0.8)
    ax.set_facecolor("white")


def save_periapsis_plot(df: pd.DataFrame) -> Path:
    out = FIG_DIR / "bmf_vs_periapsis_supported.png"
    subset = df[(df["mass_log10_kg"] == 20.0) & (df["v_inf_kms"].isin([0.0, 0.2, 0.4, 0.6, 0.8, 1.0]))].copy()
    velocities = sorted(subset["v_inf_kms"].unique())
    colors = cm.Blues(np.linspace(0.35, 0.9, len(velocities)))

    fig, ax = plt.subplots(figsize=(10.5, 5.8), dpi=220)
    support_items: list[tuple[str, int]] = []
    for velocity, color in zip(velocities, colors):
        line = subset[subset["v_inf_kms"] == velocity].sort_values("periapsis_Rm")
        support_items.append((f"v∞={velocity:.1f}", int(len(line))))
        ax.plot(
            line["periapsis_Rm"],
            line["bound_mass_fraction"],
            color=color,
            linewidth=2.2,
            marker="o",
            markersize=5.5,
            label=f"v∞ = {velocity:.1f} km/s",
            zorder=3,
        )
        ax.scatter(
            line["periapsis_Rm"],
            line["bound_mass_fraction"],
            color=color,
            edgecolor="white",
            linewidth=0.8,
            s=48,
            zorder=4,
        )

    add_axis_header(
        ax,
        "BMF falls as periapsis increases",
        f"Fixed mass = 10^20 kg | Unique physical SPH per slice: {support_note(support_items)}",
    )
    ax.set_xlabel("Periapsis (Mars radii)")
    ax.set_ylabel("Bound mass fraction")
    ax.set_ylim(-0.01, 0.32)
    ax.set_xlim(1.08, 3.02)
    style_axes(ax)
    ax.legend(frameon=False, ncol=2, fontsize=9, loc="upper right")
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def save_velocity_plot(df: pd.DataFrame) -> Path:
    out = FIG_DIR / "bmf_vs_velocity_supported.png"
    subset = df[(df["mass_log10_kg"] == 20.0) & (df["periapsis_Rm"].isin([1.2, 1.4, 1.6, 1.8, 2.0]))].copy()
    periapses = sorted(subset["periapsis_Rm"].unique())
    colors = cm.Oranges(np.linspace(0.35, 0.92, len(periapses)))

    fig, ax = plt.subplots(figsize=(10.5, 5.8), dpi=220)
    support_items: list[tuple[str, int]] = []
    for periapsis, color in zip(periapses, colors):
        line = subset[subset["periapsis_Rm"] == periapsis].sort_values("v_inf_kms")
        support_items.append((f"rp={periapsis:.1f}", int(len(line))))
        ax.plot(
            line["v_inf_kms"],
            line["bound_mass_fraction"],
            color=color,
            linewidth=2.2,
            marker="o",
            markersize=5.5,
            label=f"rp = {periapsis:.1f} R$_{{Mars}}$",
            zorder=3,
        )
        ax.scatter(
            line["v_inf_kms"],
            line["bound_mass_fraction"],
            color=color,
            edgecolor="white",
            linewidth=0.8,
            s=48,
            zorder=4,
        )

    add_axis_header(
        ax,
        "BMF also falls as velocity increases",
        f"Fixed mass = 10^20 kg | Unique physical SPH per slice: {support_note(support_items)}",
    )
    ax.set_xlabel("Encounter speed v∞ (km/s)")
    ax.set_ylabel("Bound mass fraction")
    ax.set_ylim(-0.01, 0.32)
    ax.set_xlim(-0.03, 1.03)
    style_axes(ax)
    ax.legend(frameon=False, ncol=2, fontsize=9, loc="upper right")
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def save_interaction_heatmap(df: pd.DataFrame) -> Path:
    out = FIG_DIR / "bmf_interaction_mass_velocity_periapsis.png"
    masses = [19.0, 20.0, 21.0]
    peri_vals = [1.2, 1.6, 2.0]
    vel_vals = [0.0, 0.2, 0.4, 0.6, 0.8, 1.0]
    cmap = LinearSegmentedColormap.from_list("irp_bmf", ["#eef4fb", "#7aa6dc", "#1f5aa6"])

    fig, axes = plt.subplots(1, 3, figsize=(12.0, 4.3), dpi=220, sharey=True)
    vmax = 0.28
    support_items: list[tuple[str, int]] = []
    for ax, mass in zip(axes, masses):
        panel = (
            df[(df["mass_log10_kg"] == mass) & (df["periapsis_Rm"].isin(peri_vals)) & (df["v_inf_kms"].isin(vel_vals))]
            .groupby(["periapsis_Rm", "v_inf_kms"], as_index=False)
            .agg(bound_mass_fraction=("bound_mass_fraction", "mean"), support=("physical_file", "nunique"))
        )
        support_items.append((f"10^{mass:.1f} kg", int(panel["support"].sum())))
        grid = np.full((len(peri_vals), len(vel_vals)), np.nan)
        for _, row in panel.iterrows():
            i = peri_vals.index(float(row["periapsis_Rm"]))
            j = vel_vals.index(float(row["v_inf_kms"]))
            grid[i, j] = float(row["bound_mass_fraction"])
        masked = np.ma.masked_invalid(grid)
        ax.imshow(masked, cmap=cmap, origin="lower", vmin=0.0, vmax=vmax, aspect="auto")
        ax.set_xticks(range(len(vel_vals)), [f"{v:.1f}" for v in vel_vals], fontsize=9)
        ax.set_yticks(range(len(peri_vals)), [f"{p:.1f}" for p in peri_vals], fontsize=9)
        ax.set_facecolor(GREY)
        ax.set_title(f"Mass = 10^{mass:.1f} kg", fontsize=12, color=DARK, pad=8)
        for i in range(len(peri_vals)):
            for j in range(len(vel_vals)):
                if np.isfinite(grid[i, j]):
                    ax.text(j, i, f"{grid[i, j]:.02f}", ha="center", va="center", fontsize=8, color="white")
        if ax is axes[0]:
            ax.set_ylabel("Periapsis (Mars radii)")
        ax.set_xlabel("v∞ (km/s)")
        for spine in ax.spines.values():
            spine.set_color("#aab4c0")

    fig.suptitle(
        "Mass, velocity, and periapsis interact",
        x=0.06,
        y=1.02,
        ha="left",
        fontsize=15,
        fontweight="bold",
        color=DARK,
    )
    fig.text(
        0.06,
        0.95,
        "Grey = no SPH support | Numbers are mean BMF from unique physical SPH runs in that cell",
        fontsize=9,
        color=MID_GREY,
    )
    cbar = fig.colorbar(cm.ScalarMappable(norm=Normalize(0, vmax), cmap=cmap), ax=axes, shrink=0.85, pad=0.02)
    cbar.set_label("Mean BMF", color=DARK)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def save_spin_plot(df: pd.DataFrame) -> Path:
    out = FIG_DIR / "spin_matched_comparison.png"
    subset = df[
        (df["mass_log10_kg"] == 20.0)
        & (df["periapsis_Rm"] == 1.2)
        & (df["v_inf_kms"] == 0.0)
        & (df["has_spin"])
        & (df["spin_label"].isin(["x", "y", "z", "mz"]))
    ].copy()

    fig, ax = plt.subplots(figsize=(10.5, 5.8), dpi=220)
    support_items: list[tuple[str, int]] = []
    for spin_label in ["x", "y", "z", "mz"]:
        line = subset[subset["spin_label"] == spin_label].sort_values("spin_period_hr")
        if line.empty:
            continue
        support_items.append((spin_label, int(len(line))))
        ax.plot(
            line["spin_period_hr"],
            line["bound_mass_fraction"],
            color=SPIN_COLORS[spin_label],
            linewidth=2.3,
            marker="o",
            markersize=5.5,
            label=spin_label,
        )
        ax.scatter(
            line["spin_period_hr"],
            line["bound_mass_fraction"],
            color=SPIN_COLORS[spin_label],
            edgecolor="white",
            linewidth=0.8,
            s=48,
            zorder=4,
        )

    add_axis_header(
        ax,
        "Spin can strongly change the outcome in a matched family",
        f"Fixed mass = 10^20 kg | rp = 1.2 R$_{{Mars}}$ | v∞ = 0.0 km/s | Unique physical SPH per axis: {support_note(support_items)}",
    )
    ax.set_xlabel("Spin period (hours)")
    ax.set_ylabel("Bound mass fraction")
    ax.set_xlim(2.7, 17.3)
    ax.set_ylim(-0.01, 0.31)
    style_axes(ax)
    ax.legend(title="Spin axis", frameon=False, fontsize=9, title_fontsize=10, loc="upper right")
    ax.text(
        0.02,
        0.06,
        "Same mass + same periapsis + same speed.\nOnly spin changes.",
        transform=ax.transAxes,
        fontsize=10,
        color=DARK,
        bbox={"boxstyle": "round,pad=0.35", "facecolor": "#f6f8fb", "edgecolor": "#d7dde6"},
    )
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def save_coverage_plot(df: pd.DataFrame) -> Path:
    out = FIG_DIR / "coverage_simplified.png"
    mass_vals = sorted(df["mass_log10_kg"].dropna().unique())
    peri_vals = sorted(df["periapsis_Rm"].dropna().unique())
    vel_vals = sorted(df["v_inf_kms"].dropna().unique())

    mass_peri = df.groupby(["mass_log10_kg", "periapsis_Rm"])["physical_file"].nunique().reset_index(name="count")
    peri_vel = df.groupby(["periapsis_Rm", "v_inf_kms"])["physical_file"].nunique().reset_index(name="count")

    fig, axes = plt.subplots(1, 2, figsize=(11.2, 4.8), dpi=220)
    cmap = LinearSegmentedColormap.from_list("support_blue", ["#edf4ff", "#7aa6dc", "#1f5aa6"])

    for ax, table, y_vals, x_vals, y_key, x_key, title in [
        (axes[0], mass_peri, mass_vals, peri_vals, "mass_log10_kg", "periapsis_Rm", "Mass vs periapsis"),
        (axes[1], peri_vel, peri_vals, vel_vals, "periapsis_Rm", "v_inf_kms", "Periapsis vs velocity"),
    ]:
        grid = np.full((len(y_vals), len(x_vals)), np.nan)
        for _, row in table.iterrows():
            i = y_vals.index(float(row[y_key]))
            j = x_vals.index(float(row[x_key]))
            grid[i, j] = float(row["count"])
        masked = np.ma.masked_invalid(grid)
        ax.imshow(masked, cmap=cmap, origin="lower", aspect="auto")
        ax.set_facecolor(GREY)
        ax.set_title(title, fontsize=12, color=DARK, pad=8)
        ax.set_xticks(range(len(x_vals)), [f"{v:.1f}" for v in x_vals], fontsize=8)
        ax.set_yticks(range(len(y_vals)), [f"{v:.1f}" for v in y_vals], fontsize=8)
        for i in range(len(y_vals)):
            for j in range(len(x_vals)):
                if np.isfinite(grid[i, j]):
                    ax.text(j, i, f"{int(grid[i, j])}", ha="center", va="center", fontsize=8, color="white")
        for spine in ax.spines.values():
            spine.set_color("#aab4c0")

    axes[0].set_xlabel("Periapsis (Mars radii)")
    axes[0].set_ylabel("Mass log10(kg)")
    axes[1].set_xlabel("v∞ (km/s)")
    axes[1].set_ylabel("Periapsis (Mars radii)")
    fig.suptitle("Where SPH actually exists", x=0.06, y=1.02, ha="left", fontsize=15, fontweight="bold", color=DARK)
    fig.text(0.06, 0.95, "Grey = no SPH support | Blue = unique physical SPH count", fontsize=9, color=MID_GREY)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def save_failure_table(pred_summary: pd.DataFrame) -> Path:
    out = FIG_DIR / "failure_case_table.png"
    subset = pred_summary[pred_summary["mass_log10_kg"] == 19.5].sort_values("periapsis_Rm").copy()
    subset["Periapsis"] = subset["periapsis_Rm"].map(lambda v: f"{v:.1f}")
    subset["Actual BMF"] = subset["actual_bmf"].map(lambda v: f"{v:.4f}")
    subset["OOF prediction"] = subset["oof_prediction"].map(lambda v: f"{v:.4f}")
    subset["Error"] = subset["abs_error"].map(lambda v: f"{v:.4f}")
    subset["Failure direction"] = np.where(subset["oof_prediction"] < subset["actual_bmf"], "under", "over")
    display = subset[["Periapsis", "Actual BMF", "OOF prediction", "Error", "Failure direction"]]

    fig, ax = plt.subplots(figsize=(10.5, 3.7), dpi=220)
    ax.axis("off")
    ax.set_title(
        "Mass 10^19.5 kg is in range but not well supported",
        loc="left",
        fontsize=15,
        fontweight="bold",
        color=DARK,
        pad=16,
    )
    ax.text(
        0.0,
        0.96,
        "Only 2 unique physical SPH cases support this slice. The held-out surrogate misses in opposite directions.",
        transform=ax.transAxes,
        fontsize=9.5,
        color=MID_GREY,
    )
    table = ax.table(
        cellText=display.values,
        colLabels=list(display.columns),
        loc="center",
        cellLoc="center",
        colLoc="center",
    )
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1, 1.65)
    for (row, col), cell in table.get_celld().items():
        cell.set_edgecolor("#d8dee8")
        if row == 0:
            cell.set_facecolor(BLUE)
            cell.get_text().set_color("white")
            cell.get_text().set_weight("bold")
        else:
            cell.set_facecolor("#f7f9fc" if row % 2 else "white")
            if col == 4:
                cell.get_text().set_color(ACCENT if "over" in cell.get_text().get_text() else BLUE)
                cell.get_text().set_weight("bold")
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def save_schematic() -> Path:
    out = FIG_DIR / "tidal_disruption_schematic.png"
    fig, ax = plt.subplots(figsize=(6.6, 4.4), dpi=220)
    ax.set_aspect("equal")
    ax.axis("off")
    mars = plt.Circle((0, 0), 1.0, color="#c84d43", alpha=0.95)
    tidal = plt.Circle((0, 0), 1.85, fill=False, linestyle="--", linewidth=2.0, color=LIGHT_BLUE)
    ax.add_patch(mars)
    ax.add_patch(tidal)
    t = np.linspace(-1.6, 1.4, 200)
    x = 3.1 - 0.9 * (t + 0.2) ** 2
    y = 1.6 * t
    ax.plot(x, y, color=DARK, linewidth=2.2)
    ax.arrow(x[-15], y[-15], 0.001, 0.001, head_width=0.12, head_length=0.18, fc=DARK, ec=DARK)
    fragment_x = np.array([1.15, 1.28, 1.42, 1.55, 1.72])
    fragment_y = np.array([0.42, 0.16, -0.05, -0.18, -0.33])
    ax.scatter(fragment_x, fragment_y, s=[120, 80, 58, 42, 32], color=ACCENT, edgecolor="white", linewidth=0.8)
    ax.text(-0.1, 0.0, "Mars", color="white", fontsize=14, fontweight="bold", ha="center", va="center")
    ax.text(0.0, 2.0, "strong tidal region", color=LIGHT_BLUE, fontsize=10, ha="center")
    ax.text(2.3, 1.9, "incoming asteroid", color=DARK, fontsize=10)
    ax.text(1.9, -0.65, "bound + escaping fragments", color=ACCENT, fontsize=10)
    ax.text(-1.9, -2.2, "schematic, not to scale", fontsize=9.5, color=MID_GREY)
    ax.set_xlim(-2.4, 3.2)
    ax.set_ylim(-2.4, 2.5)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight", transparent=True)
    plt.close(fig)
    return out


def save_dashboard_mock() -> Path:
    out = FIG_DIR / "screening_dashboard_mock.png"
    fig, ax = plt.subplots(figsize=(10.8, 5.8), dpi=220)
    ax.axis("off")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)

    ax.add_patch(plt.Rectangle((0.3, 0.3), 9.4, 5.4, facecolor=PALE, edgecolor="#dbe3ef", linewidth=1.2))
    ax.text(0.6, 5.25, "Screening dashboard", fontsize=18, fontweight="bold", color=DARK)
    ax.text(0.6, 4.92, "Use ML as a screening tool first. Escalate to SPH where support is weak.", fontsize=10, color=MID_GREY)

    ax.add_patch(plt.Rectangle((0.6, 1.0), 3.9, 3.5, facecolor="white", edgecolor="#dbe3ef", linewidth=1.1))
    ax.text(0.85, 4.1, "Inputs", fontsize=13, fontweight="bold", color=DARK)
    ax.text(0.85, 3.65, "mass / size", fontsize=11, color=DARK)
    ax.text(0.85, 3.25, "periapsis", fontsize=11, color=DARK)
    ax.text(0.85, 2.85, "eccentricity", fontsize=11, color=DARK)
    ax.text(0.85, 2.45, "spin axis + period", fontsize=11, color=DARK)
    ax.text(0.85, 1.88, "Rocky density assumed = 2700 kg/m^3", fontsize=11, color=BLUE, fontweight="bold")

    ax.add_patch(plt.Rectangle((5.0, 1.0), 4.1, 3.5, facecolor="white", edgecolor="#dbe3ef", linewidth=1.1))
    ax.text(5.25, 4.1, "Outputs", fontsize=13, fontweight="bold", color=DARK)
    ax.text(5.25, 3.65, "predicted BMF", fontsize=11, color=DARK)
    ax.text(5.25, 3.25, "bound mass", fontsize=11, color=DARK)
    ax.text(5.25, 2.85, "support level", fontsize=11, color=DARK)
    ax.text(5.25, 2.45, "SPH required? yes / no", fontsize=11, color=DARK)
    ax.text(5.25, 1.88, "screening threshold = BMF >= 10%", fontsize=11, color=ACCENT, fontweight="bold")

    ax.annotate("", xy=(4.9, 2.7), xytext=(4.55, 2.7), arrowprops=dict(arrowstyle="->", color=DARK, lw=2))
    ax.text(6.65, 0.62, "Low support, edge cases, and borderline BMF go back to SPH.", fontsize=10, color=MID_GREY, ha="center")
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
    run.font.bold = bold
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(8.6), Inches(0.7), title, font_size=26, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(8.9), Inches(0.35), subtitle, font_size=11, color=MID_GREY)


def add_full_image(slide, path: Path, left=0.5, top=1.25, width=12.33, height=5.45) -> None:
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_footer(slide, page_number: int) -> None:
    add_textbox(slide, Inches(12.55), Inches(7.0), Inches(0.5), Inches(0.25), str(page_number), font_size=9, color=MID_GREY, align=PP_ALIGN.RIGHT)


def build_presentation(figures: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Mars Tidal Encounters: Physics First", "IRP refactor | 10 slides | August 11, 2026")
    add_textbox(
        slide,
        Inches(0.5),
        Inches(1.15),
        Inches(5.2),
        Inches(1.2),
        "Main question: when do close asteroid flybys leave Mars-bound debris?",
        font_size=22,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(2.25),
        Inches(4.8),
        Inches(1.8),
        "This version leads with SPH-supported physical findings.\nML is used only as an in-domain screening layer.",
        font_size=16,
    )
    slide.shapes.add_picture(str(figures["schematic"]), Inches(6.1), Inches(1.35), width=Inches(6.5), height=Inches(5.2))
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Four Results", "Physics findings first. Screening logic second.")
    items = [
        "1. Closer and slower encounters retain more bound mass.",
        "2. Mass, velocity, and periapsis interact rather than acting independently.",
        "3. Spin can strongly change BMF in some matched regimes.",
        "4. The surrogate is trusted only where local SPH support exists.",
    ]
    top = 1.45
    for i, item in enumerate(items):
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top + i * 1.2), Inches(11.7), Inches(0.72))
        shape = slide.shapes[-1]
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string("F4F7FB")
        shape.line.color.rgb = RGBColor.from_string("D7DFEA")
        add_textbox(slide, Inches(1.05), Inches(top + 0.12 + i * 1.2), Inches(11.0), Inches(0.4), item, font_size=18)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.45),
        Inches(9.5),
        Inches(0.45),
        "Archive used here: 279 unique physical SPH runs after removing repeated FoF rows.",
        font_size=11,
        color=MID_GREY,
    )
    add_footer(slide, 2)

    # Slides 3-9 with figures
    slide_specs = [
        ("Closer passes keep more bound mass", "Controlled SPH slice: fixed mass, multiple velocities.", figures["periapsis"]),
        ("Faster flybys suppress retention", "Controlled SPH slice: fixed mass, multiple periapses.", figures["velocity"]),
        ("The three-way interaction matters", "Actual SPH means only. Grey cells were never simulated.", figures["interaction"]),
        ("Spin can flip the outcome", "Matched family: same mass, same periapsis, same speed.", figures["spin"]),
        ("Trust depends on local SPH support", "Blue = unique physical SPH count. Grey = no support.", figures["coverage"]),
        ("A concrete failure case", "Mass 10^19.5 kg looks in-range globally but fails locally.", figures["failure"]),
        ("How the screening tool should be presented", "Inputs, outputs, and wording aligned to the physics-first story.", figures["dashboard"]),
    ]
    for idx, (title, subtitle, fig_path) in enumerate(slide_specs, start=3):
        slide = prs.slides.add_slide(blank)
        add_title(slide, title, subtitle)
        add_full_image(slide, fig_path)
        add_footer(slide, idx)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Takeaways", "What should carry into the spoken summary.")
    items = [
        "Closer and slower flybys are the clearest route to higher BMF.",
        "Mass lifts retention, but only together with favourable periapsis and velocity.",
        "Spin is not a small correction everywhere; some matched SPH families change sharply with spin.",
        "Use ML for screening, not replacement. Sparse or weakly supported regions still need SPH.",
    ]
    top = 1.45
    for i, item in enumerate(items):
        add_textbox(slide, Inches(0.9), Inches(top + i * 1.15), Inches(11.0), Inches(0.55), item, font_size=18)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(6.25),
        Inches(11.3),
        Inches(0.6),
        "Immediate SPH priority: fill the sparse mass 10^19.5 kg region and other grey cells before widening surrogate claims.",
        font_size=12,
        color=MID_GREY,
    )
    add_footer(slide, 10)

    prs.save(OUT_PATH)


def main() -> None:
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    frame = load_bound_frame()
    unique_runs = dedupe_physical_runs(frame)
    pred_summary = load_prediction_summary()
    figures = {
        "schematic": save_schematic(),
        "periapsis": save_periapsis_plot(unique_runs),
        "velocity": save_velocity_plot(unique_runs),
        "interaction": save_interaction_heatmap(unique_runs),
        "spin": save_spin_plot(unique_runs),
        "coverage": save_coverage_plot(unique_runs),
        "failure": save_failure_table(pred_summary),
        "dashboard": save_dashboard_mock(),
    }
    build_presentation(figures)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
